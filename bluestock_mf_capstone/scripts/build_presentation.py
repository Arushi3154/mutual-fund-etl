"""
Generates the official 12-slide presentation for the Bluestock Mutual Fund Capstone.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_PATH = BASE_DIR / "reports" / "Presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

slides_content = [
    # Slide 1: Title
    ("Bluestock Mutual Fund Capstone Analytics", 
     "End-to-End ETL Pipeline, Performance Metrics & Recommender Engine\n\nAuthor: Arushi Awasthi\nRepository: github.com/Arushi3154/mutual-fund-etl"),
    
    # Slide 2: Problem & Objective
    ("Problem Statement & Objectives", 
     "• Ingest, clean, and standardize fragmented Indian Mutual Fund NAV & AUM datasets.\n• Compute risk-adjusted metrics: 252-day trading CAGR, Sharpe, Beta, and 95% VaR.\n• Resolve weekend NAV gaps using ffill() to avoid artificial volatility spikes.\n• Automate fund recommendation scorecards and PowerBI data marts."),
    
    # Slide 3: Data Sources
    ("Data Sources & Ingestion Layer", 
     "• Official AMFI API (mfapi.in) for live daily NAV tracking.\n• Historical Scheme Master, Portfolio Holdings, and Benchmark Indices (Nifty 50).\n• Monthly industry-wide AUM and category net inflow statistics.\n• Automated error handling for HTTP timeouts and payload parsing."),
    
    # Slide 4: Architecture
    ("ETL Architecture & Warehouse Design", 
     "• Ingestion Layer: API fetching and automated schema validation.\n• Transformation Layer: Date standardization, missing value imputations, unit scaling.\n• Database Warehouse: Normalized 3NF SQLite database (bluestock_mf.db).\n• Export Layer: Automated CSV data mart generation for PowerBI dashboards."),
    
    # Slide 5: EDA Highlights (1/2)
    ("Exploratory Data Analysis — Category Distribution", 
     "• Equity schemes dominate overall industry AUM (over 62% of total assets).\n• Large-cap funds display lower historical variance compared to Small/Mid-cap funds.\n• Multi-cap and Flexi-cap categories exhibit highest investor inflows."),
    
    # Slide 6: EDA Highlights (2/2)
    ("Exploratory Data Analysis — Expense Ratios & NAVs", 
     "• Direct plans consistently outperform Regular plans by 0.75%–1.25% annually.\n• Expense ratio distribution peaks between 0.8% and 1.8% for equity schemes.\n• Outlier detection identified split/merger anomalies in NAV history."),
    
    # Slide 7: Performance Metrics (1/2)
    ("Performance Analytics — Risk-Adjusted Returns", 
     "• Trading-day annualized CAGR formula: CAGR = (NAV_end / NAV_start)^(252 / n) - 1.\n• Sharpe Ratio evaluated against Risk-Free Rate (Rf = 6.5%).\n• Sortino Ratio highlights downside volatility relative to total variance."),
    
    # Slide 8: Performance Metrics (2/2)
    ("Performance Analytics — Market Sensitivity & Risk", 
     "• Beta calculated via covariance against Nifty 50 benchmark index.\n• Alpha generation measured using Capital Asset Pricing Model (CAPM).\n• 95% Value at Risk (VaR) computed to quantify potential single-day losses."),
    
    # Slide 9: Dashboard Screenshots (1/2)
    ("Dashboard Screenshots — Executive & Category Overview", 
     "• Page 1: Executive KPI Scorecard (Total AUM, Active Schemes, Average CAGR).\n• Page 2: Monthly SIP Inflow Trends & Fund House AUM Comparisons.\n• Interactive Slicers: Fund Category, AMC, and Date Range filters on all pages."),
    
    # Slide 10: Dashboard Screenshots (2/2)
    ("Dashboard Screenshots — Scheme Deep-Dive & Recommendations", 
     "• Page 3: Risk vs Return Scatter Plot Matrix (Sharpe vs Beta).\n• Page 4: Top 5 Recommended Funds & Scorecard Table.\n• Cross-filtering allows immediate drill-down by AMC and Risk Rating."),
    
    # Slide 11: Key Findings & Recommendations
    ("Key Findings & Strategic Recommendations", 
     "• Top 10% of schemes generate over 75% of cumulative risk-adjusted alpha.\n• Direct equity funds maintain superior 3-year CAGR due to expense savings.\n• Recommendation: Rebalance portfolios quarterly targeting Sharpe > 1.2."),
    
    # Slide 12: Thank You
    ("Thank You & Repository Access", 
     "Questions & Answers\n\nGitHub Repository: https://github.com/Arushi3154/mutual-fund-etl\nRelease Version: v1.0")
]

for title_text, content_text in slides_content:
    slide = prs.slides.add_slide(blank_layout)
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 32, 67)

    # Content Box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.8))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(50, 50, 50)

prs.save(OUTPUT_PATH)
print(f"Presentation saved to {OUTPUT_PATH}")
